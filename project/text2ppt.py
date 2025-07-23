from fastapi import FastAPI, Form, Request, UploadFile, File, HTTPException
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse, RedirectResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from datetime import datetime
from typing import List, Optional
import os
import json
import shutil
import re
from collections import Counter
import PyPDF2
import matplotlib.pyplot as plt
import pandas as pd
import uuid
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
MAX_TEXT_WORDS = 1500  # Limit input text to 1500 words
MAX_IMAGE_SIZE = 2 * 1024 * 1024  # 2 MB per image


try:
    import nltk
    from nltk.tokenize import sent_tokenize, word_tokenize
    from nltk.corpus import stopwords
    from nltk.stem import WordNetLemmatizer
    NLTK_AVAILABLE = True
    print("NLTK imported successfully")
except ImportError as e:
    print(f"NLTK import failed: {e}")
    NLTK_AVAILABLE = False


class EnhancedSummarizer:
    def __init__(self):
        self.nltk_available = NLTK_AVAILABLE
        
        if self.nltk_available:
            try:
                
                self._download_nltk_data()
                self.stop_words = set(stopwords.words('english'))
                self.lemmatizer = WordNetLemmatizer()
                print("NLTK initialized successfully")
            except Exception as e:
                print(f"NLTK initialization failed: {e}")
                self.nltk_available = False
                self._setup_fallback()
        else:
            self._setup_fallback()
    
    def _download_nltk_data(self):
        """Download required NLTK data"""
        import ssl
        try:
            _create_unverified_https_context = ssl._create_unverified_context
        except AttributeError:
            pass
        else:
            ssl._create_default_https_context = _create_unverified_https_context
        
        
        datasets = ['punkt', 'stopwords', 'wordnet', 'omw-1.4']
        for dataset in datasets:
            try:
                nltk.data.find(f'tokenizers/{dataset}' if dataset == 'punkt' else f'corpora/{dataset}')
                print(f"NLTK {dataset} already available")
            except LookupError:
                try:
                    print(f"Downloading NLTK {dataset}...")
                    nltk.download(dataset, quiet=True)
                    print(f"NLTK {dataset} downloaded successfully")
                except Exception as e:
                    print(f"Failed to download {dataset}: {e}")
    
    def _setup_fallback(self):
        """Setup fallback stopwords and methods"""
        print("Setting up fallback text processing")
        self.stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by',
            'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did',
            'will', 'would', 'should', 'could', 'may', 'might', 'must', 'can', 'this', 'that',
            'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they', 'me', 'him', 'her',
            'us', 'them', 'my', 'your', 'his', 'her', 'its', 'our', 'their', 'what', 'which',
            'who', 'whom', 'am', 'being', 'about', 'above', 'after', 'again', 'against', 'all',
            'any', 'because', 'before', 'below', 'between', 'both', 'down', 'during', 'each',
            'few', 'from', 'further', 'here', 'how', 'into', 'more', 'most', 'no', 'nor', 'not',
            'now', 'off', 'once', 'only', 'other', 'out', 'over', 'own', 'same', 'so', 'some',
            'such', 'than', 'then', 'there', 'through', 'too', 'under', 'until', 'up', 'very',
            'where', 'while', 'why', 'with'
        }
    
    def clean_text(self, text):
        """Clean and preprocess text"""
        
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s.,!?;:]', '', text)
        return text.strip()
    
    def simple_sentence_split(self, text):
        """Fallback sentence splitting without NLTK"""
        
        sentences = re.split(r'[.!?]+', text)
        return [s.strip() for s in sentences if s.strip()]
    
    def simple_word_tokenize(self, text):
        """Fallback word tokenization without NLTK"""
        
        words = re.findall(r'\b\w+\b', text.lower())
        return words
    
    def extract_key_phrases(self, text, max_phrases=10):
        """Extract key phrases from text"""
        try:
            if self.nltk_available:
                sentences = sent_tokenize(text)
            else:
                sentences = self.simple_sentence_split(text)
        except:
            sentences = self.simple_sentence_split(text)
        
        
        key_phrases = []
        for sentence in sentences:
            try:
                if self.nltk_available:
                    words = word_tokenize(sentence.lower())
                else:
                    words = self.simple_word_tokenize(sentence)
            except:
                words = self.simple_word_tokenize(sentence)
            
            
            filtered_words = [word for word in words if word not in self.stop_words and len(word) > 2]
            
            
            for i in range(len(filtered_words) - 1):
                phrase = ' '.join(filtered_words[i:i+2])
                key_phrases.append(phrase)
                
            
            key_phrases.extend(filtered_words)
        
        
        phrase_counts = Counter(key_phrases)
        return [phrase for phrase, count in phrase_counts.most_common(max_phrases)]
    
    def generate_bullet_points(self, text, max_bullets=3):
        """Generate bullet points from text using extractive summarization"""
        try:
            if self.nltk_available:
                sentences = sent_tokenize(text)
            else:
                sentences = self.simple_sentence_split(text)
        except:
            sentences = self.simple_sentence_split(text)
        
        if len(sentences) <= max_bullets:
            return [self.clean_sentence(sent) for sent in sentences if len(sent.strip()) > 10]
        
        
        key_phrases = self.extract_key_phrases(text)
        sentence_scores = {}
        
        for i, sentence in enumerate(sentences):
            score = 0
            sentence_lower = sentence.lower()
            
            
            for phrase in key_phrases:
                if phrase in sentence_lower:
                    score += 1
            
            
            if re.search(r'\d+', sentence):
                score += 0.5
            if any(keyword in sentence_lower for keyword in ['important', 'key', 'main', 'significant', 'crucial']):
                score += 0.5
            
            
            if i < len(sentences) * 0.3:
                score += 0.3
            
            sentence_scores[i] = score
        
        
        top_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)[:max_bullets]
        top_sentences.sort(key=lambda x: x[0])  
        
        bullets = []
        for idx, score in top_sentences:
            bullet = self.clean_sentence(sentences[idx])
            if bullet and len(bullet) > 10:
                bullets.append(bullet)
        
        return bullets[:max_bullets]
    
    def clean_sentence(self, sentence):
        """Clean and format a sentence for bullet points"""
        
        sentence = re.sub(r'\s+', ' ', sentence).strip()
        
        sentence = re.sub(r'^(and|but|or|so|because|however|therefore|thus|hence)\s+', '', sentence, flags=re.IGNORECASE)
        
        if sentence:
            sentence = sentence[0].upper() + sentence[1:]
        return sentence
    
    def generate_title(self, text, max_words=4):
        """Generate a meaningful title from text"""
        try:
            
            key_phrases = self.extract_key_phrases(text, max_phrases=5)
            
            if not key_phrases:
                return "Key Points"
            
            
            title_words = []
            for phrase in key_phrases:
                words = phrase.split()
                for word in words:
                    if len(word) > 2 and word not in self.stop_words:
                        title_words.append(word.capitalize())
                        if len(title_words) >= max_words:
                            break
                if len(title_words) >= max_words:
                    break
            
            if not title_words:
                
                try:
                    if self.nltk_available:
                        words = word_tokenize(text.lower())
                    else:
                        words = self.simple_word_tokenize(text)
                except:
                    words = self.simple_word_tokenize(text)
                
                title_words = [word.capitalize() for word in words[:max_words*2] 
                              if word not in self.stop_words and len(word) > 2][:max_words]
            
            return ' '.join(title_words[:max_words]) if title_words else "Overview"
        except Exception as e:
            print(f"Error generating title: {e}")
            return "Overview"


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
PPT_FOLDER = os.path.join(BASE_DIR, "generated_ppt")
LOG_FILE = os.path.join(BASE_DIR, "logs.json")
USER_FILE = os.path.join(BASE_DIR, "users.json")

os.makedirs(TEMPLATES_DIR, exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PPT_FOLDER, exist_ok=True)

if not os.path.exists(USER_FILE):
    with open(USER_FILE, "w") as f:
        json.dump({}, f)


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"]
)


summarizer = EnhancedSummarizer()


def load_users():
    with open(USER_FILE, "r") as f:
        return json.load(f)

def save_users(users):
    with open(USER_FILE, "w") as f:
        json.dump(users, f, indent=2)

def load_logs():
    if not os.path.exists(LOG_FILE):
        return {}
    with open(LOG_FILE, "r") as f:
        return json.load(f)

def save_logs(logs):
    with open(LOG_FILE, "w") as f:
        json.dump(logs, f, indent=4)

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting PDF: {e}")
    return text

def add_slide_design(slide):
    """Add professional design elements to slide"""
    try:
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250) 
    except:
        pass

def add_fade_transition(slide):
    """Add fade transition to slide"""
    try:
        slide._element.insert(2, parse_xml("""
            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                <p:fade />
            </p:transition>
        """))
    except:
        pass 

def add_header_footer(slide, username, page_number, total_pages):
    """Add header and footer to slide"""
    try:
        
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
        header_frame = header.text_frame
        header_frame.text = "Auto-Generated Presentation"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(14)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(70, 130, 180)  # Steel blue
        header_para.alignment = PP_ALIGN.CENTER
        
        
        footer = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(1.5), Inches(0.4))
        footer_frame = footer.text_frame
        footer_frame.text = f"{page_number}/{total_pages}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(100, 100, 100)
        footer_para.alignment = PP_ALIGN.RIGHT
    except Exception as e:
        print(f"Error adding header/footer: {e}")

def create_title_slide(prs, title, subtitle, username):
    """Create a professional title slide"""
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    
    
    add_slide_design(slide)
    
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_frame = title_placeholder.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(25, 25, 112) 
    title_para.alignment = PP_ALIGN.CENTER
    
    
    try:
        if len(slide.shapes.placeholders) > 1:
            subtitle_placeholder = slide.shapes.placeholders[1]
            subtitle_placeholder.text = subtitle
            subtitle_frame = subtitle_placeholder.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(70, 130, 180)
            subtitle_para.alignment = PP_ALIGN.CENTER
    except:
        pass
    
    
    try:
        author_box = slide.shapes.add_textbox(Inches(7.5), Inches(6.5), Inches(2.5), Inches(0.8))
        author_frame = author_box.text_frame
        author_frame.text = f"Generated by: {username}\n{datetime.now().strftime('%B %d, %Y')}"
        author_para = author_frame.paragraphs[0]
        author_para.font.size = Pt(12)
        author_para.font.color.rgb = RGBColor(100, 100, 100)
        author_para.alignment = PP_ALIGN.RIGHT
    except:
        pass
    
    add_fade_transition(slide)
    return slide

def create_content_slide(prs, title, bullets, page_number, total_pages, username, image_path=None):
    """Create a content slide with bullets"""
    bullet_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_layout)
    
    
    add_slide_design(slide)
    
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_frame = title_placeholder.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(25, 25, 112)
    
    
    content_placeholder = slide.shapes.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            
            p = text_frame.paragraphs[0]
        else:
            
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_after = Pt(12)
    
    
    if image_path and os.path.exists(image_path):
        try:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            image_width = Inches(2.0)
            image_height = Inches(1.5)
            
            
            left = slide_width - image_width - Inches(0.5)
            top = slide_height - image_height - Inches(0.5)
            
            slide.shapes.add_picture(image_path, left, top, image_width, image_height)
        except Exception as e:
            print(f"Error adding image: {e}")
    
    add_header_footer(slide, username, page_number, total_pages)
    add_fade_transition(slide)
    return slide

def create_enhanced_ppt(text, reference, username, image_paths=None):
    """Create an enhanced PowerPoint presentation"""
    try:
        prs = Presentation()
        image_paths = image_paths or []
        
        
        text = summarizer.clean_text(text)
        
        
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if not paragraphs:
            
            try:
                if summarizer.nltk_available:
                    sentences = sent_tokenize(text)
                else:
                    sentences = summarizer.simple_sentence_split(text)
            except:
                sentences = summarizer.simple_sentence_split(text)
            
            chunk_size = max(3, len(sentences) // 8)  
            paragraphs = []
            for i in range(0, len(sentences), chunk_size):
                chunk = ' '.join(sentences[i:i+chunk_size])
                if len(chunk.strip()) > 50:  
                    paragraphs.append(chunk)
        
        
        paragraphs = paragraphs[:10]  
        
        
        total_slides = len(paragraphs) + 2  
        if reference.strip():
            total_slides += 1
        
        
        main_title = summarizer.generate_title(text, max_words=6)
        subtitle = "AI-Generated Presentation"
        create_title_slide(prs, main_title, subtitle, username)
        
        
        slide_count = 1
        for idx, paragraph in enumerate(paragraphs):
            if len(paragraph.strip()) < 20:  
                continue
                
            
            slide_title = summarizer.generate_title(paragraph, max_words=4)
            bullets = summarizer.generate_bullet_points(paragraph, max_bullets=3)
            
            if bullets:  
                slide_count += 1
                image_path = image_paths[idx] if idx < len(image_paths) else None
                create_content_slide(prs, slide_title, bullets, slide_count, total_slides, username, image_path)
        
        
        if reference.strip():
            slide_count += 1
            ref_bullets = [ref.strip() for ref in reference.split('\n') if ref.strip()]
            if not ref_bullets:
                ref_bullets = [reference.strip()]
            create_content_slide(prs, "References", ref_bullets, slide_count, total_slides, username)
        
       
        slide_count += 1
        try:
            thank_you_layout = prs.slide_layouts[6]  
        except:
            thank_you_layout = prs.slide_layouts[0]  
        
        thank_slide = prs.slides.add_slide(thank_you_layout)
        add_slide_design(thank_slide)
        
        thank_box = thank_slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        thank_frame = thank_box.text_frame
        thank_frame.text = "Thank You!"
        thank_para = thank_frame.paragraphs[0]
        thank_para.font.size = Pt(48)
        thank_para.font.bold = True
        thank_para.font.color.rgb = RGBColor(25, 25, 112)
        thank_para.alignment = PP_ALIGN.CENTER
        
        
        contact_box = thank_slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
        contact_frame = contact_box.text_frame
        contact_frame.text = "Questions & Discussion"
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(24)
        contact_para.font.color.rgb = RGBColor(70, 130, 180)
        contact_para.alignment = PP_ALIGN.CENTER
        
        add_header_footer(thank_slide, username, slide_count, total_slides)
        add_fade_transition(thank_slide)
        

        filename = f"{username}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(PPT_FOLDER, filename)
        prs.save(filepath)
        return filepath
        
    except Exception as e:
        print(f"Error in create_enhanced_ppt: {e}")
        import traceback
        traceback.print_exc()
        raise e

# --- Routes ---
@app.get("/")
async def home():
    return RedirectResponse("/signinsignup.html")

@app.get("/signinsignup.html")
async def signin_signup():
    try:
        with open(os.path.join(TEMPLATES_DIR, "signinsignup.html"), "r", encoding="utf-8") as f:
            return HTMLResponse(f.read())
    except:
        return HTMLResponse("<h1>signinsignup.html not found</h1>", status_code=500)

@app.get("/main.html")
async def main_page():
    try:
        with open(os.path.join(TEMPLATES_DIR, "main.html"), "r", encoding="utf-8") as f:
            return HTMLResponse(f.read())
    except:
        return HTMLResponse("<h1>main.html not found</h1>", status_code=500)

@app.post("/signup")
async def signup(username: str = Form(...), name: str = Form(...), password: str = Form(...), email: str = Form(...)):
    users = load_users()
    if username in users:
        return HTMLResponse("<h2>User already exists</h2>", status_code=400)
    users[username] = {"name": name, "password": password, "email": email}
    save_users(users)
    return RedirectResponse("/signinsignup.html", status_code=303)

@app.post("/signin")
async def signin(username: str = Form(...), password: str = Form(...)):
    users = load_users()
    if username in users and users[username]["password"] == password:
        return RedirectResponse(f"/main.html?user={username}", status_code=303)
    return HTMLResponse("<h2>Invalid login</h2>", status_code=401)

def generate_graph_from_csv(file_path, username):
        try:
            df = pd.read_csv(file_path)
            print(f"CSV Columns: {df.columns.tolist()}")

            if df.empty or len(df.columns) < 2:
                print("CSV does not have enough data for graphing.")
                return None

            plt.figure(figsize=(8, 5))
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) < 2:
                print("Need at least two numeric columns for a meaningful graph.")
                return None

            x = numeric_cols[0]
            for y in numeric_cols[1:]:
                plt.plot(df[x], df[y], label=f"{y}")

            plt.xlabel(x)
            plt.title("CSV Data Graph")
            plt.legend()
            plt.grid(True)

            graph_path = os.path.join(UPLOAD_FOLDER, f"{username}_graph.png")
            plt.tight_layout()
            plt.savefig(graph_path)
            plt.close()
            print(f"Graph saved to {graph_path}")
            return graph_path
        except Exception as e:
            print(f"Failed to generate graph: {e}")
            return None

@app.post("/generate-ppt")
async def generate_ppt(
    request: Request,
    text: str = Form(""),
    reference: str = Form(""),
    images: list[UploadFile] = File(default=[]),
    doc: UploadFile = File(None),
    csv_file: UploadFile = File(None),
    username: str = Form(...)
):
    ip = request.client.host or "unknown"
    browser = request.headers.get("user-agent", "unknown")
    logs = load_logs()

    
    if doc and doc.filename:
        contents = await doc.read()
        if len(contents) > MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail="Uploaded document is too large (maximum 10 MB allowed)")
        
        
        doc.file.seek(0)

        doc_path = os.path.join(UPLOAD_FOLDER, f"{username}_{doc.filename}")
        with open(doc_path, "wb") as buffer:
            shutil.copyfileobj(doc.file, buffer)
        
        try:
            if doc.filename.endswith(".docx"):
                document = Document(doc_path)
                text = "\n".join([p.text for p in document.paragraphs if p.text.strip()])
            elif doc.filename.endswith(".txt"):
                with open(doc_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif doc.filename.endswith(".pdf"):
                text = extract_text_from_pdf(doc_path)
            
            
            word_count = len(text.split())
            if word_count > MAX_TEXT_WORDS:
                raise HTTPException(
                    status_code=400,
                    detail=f"Text extracted from the document is too long ({word_count} words). Limit is {MAX_TEXT_WORDS} words."
                )
        finally:
            
            if os.path.exists(doc_path):
                os.remove(doc_path)


    
    if not text.strip():
        return HTMLResponse("<h2>Please provide text content or upload a document</h2>", status_code=400)

    
    logs[username] = {
        "ip": ip,
        "browser": browser,
        "history": logs.get(username, {}).get("history", []) + [text.strip()[:100]]  # Store first 100 chars
    }
    save_logs(logs)

    
    image_paths = []
    for image in images:
        if image.filename:
           
            contents = await image.read()
            if len(contents) > MAX_IMAGE_SIZE:
                raise HTTPException(
                    status_code=413,
                    detail=f"Image '{image.filename}' is too large (maximum 2 MB allowed)"
                )

            
            image.file.seek(0)

            image_path = os.path.join(UPLOAD_FOLDER, f"{datetime.now().timestamp()}_{image.filename}")
            with open(image_path, "wb") as buffer:
                shutil.copyfileobj(image.file, buffer)
            image_paths.append(image_path)

    try:
        
        ppt_path = create_enhanced_ppt(text, reference, username, image_paths)
        
        
        for image_path in image_paths:
            if os.path.exists(image_path):
                os.remove(image_path)
        
        return FileResponse(
            ppt_path, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            filename="enhanced_presentation.pptx"
        )
    except Exception as e:
        
        for image_path in image_paths:
            if os.path.exists(image_path):
                os.remove(image_path)
         # Handle CSV and generate graph
        graph_image_path = None
        if csv_file and csv_file.filename.endswith(".csv"):
            csv_path = os.path.join(UPLOAD_FOLDER, f"{username}_{csv_file.filename}")
            with open(csv_path, "wb") as buffer:
                shutil.copyfileobj(csv_file.file, buffer)
            graph_image_path = generate_graph_from_csv(csv_path, username)
            if os.path.exists(csv_path):
                os.remove(csv_path)
            if graph_image_path:
                image_paths.append(graph_image_path)

        # Generate PowerPoint
        try:
            ppt_path = create_enhanced_ppt(text, reference, username, image_paths)
            for path in image_paths:
                if os.path.exists(path):
                    os.remove(path)
            return FileResponse(
                ppt_path,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename="enhanced_presentation.pptx"
            )
        except Exception as e:
            for path in image_paths:
                if os.path.exists(path):
                    os.remove(path)
            raise HTTPException(status_code=500, detail=f"Error generating PPT: {str(e)}")

@app.post("/generate-from-doc")
async def generate_from_doc(
    request: Request,
    docFile: UploadFile = File(...),
    username: str = Form(...),
    reference: str = Form(None)
):
    ip = request.client.host or "unknown"
    browser = request.headers.get("user-agent", "unknown")
    logs = load_logs()

    path = os.path.join(UPLOAD_FOLDER, f"{username}_{uuid.uuid4()}_{docFile.filename}")
    extension = os.path.splitext(docFile.filename)[-1].lower()

    try:
        with open(path, "wb") as buffer:
            shutil.copyfileobj(docFile.file, buffer)
        print(f"Uploaded file saved at: {path}")

        # CASE 1: CSV FILE
        if extension == ".csv":
            try:
                try:
                    df = pd.read_csv(path)  # Try UTF-8
                except UnicodeDecodeError:
                    df = pd.read_csv(path, encoding='latin1')  # Fallback

                if df.shape[1] < 2:
                    raise HTTPException(status_code=400, detail="CSV must have at least 2 columns.")
                
                prs = Presentation()
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                title_slide.shapes.title.text = f"CSV Visualization by {username}"
                title_slide.placeholders[1].text = "Auto-generated from CSV"

                for i in range(1, df.shape[1]):
                    chart_data = CategoryChartData()
                    chart_data.categories = list(df.iloc[:, 0])
                    chart_data.add_series(df.columns[i], list(df.iloc[:, i]))

                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = f"{df.columns[i]} Chart"
                    slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(1), Inches(1.5), Inches(8), Inches(4.5),
                        chart_data
                    )
                
                if reference:
                    ref_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    ref_slide.shapes.title.text = "Reference"
                    ref_slide.placeholders[1].text = reference

                output_path = f"csv_generated_{uuid.uuid4().hex[:6]}.pptx"
                prs.save(output_path)

                return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="CSV_Generated_Slides.pptx")

            except Exception as e:
                raise HTTPException(status_code=400, detail=f"CSV Error: {str(e)}")


        # CASE 2: DOCX, TXT, PDF
        content = ""
        if extension == ".docx":
            document = Document(path)
            content = "\n".join([p.text for p in document.paragraphs if p.text.strip()])
        elif extension == ".txt":
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
        elif extension == ".pdf":
            content = extract_text_from_pdf(path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format. Only .csv, .docx, .txt, and .pdf are allowed.")

        if not content.strip():
            raise HTTPException(status_code=400, detail="No text content found in the document")

        logs[username] = {
            "ip": ip,
            "browser": browser,
            "history": logs.get(username, {}).get("history", []) + [content.strip()[:100]]
        }
        save_logs(logs)

        print("Generating PPT from document text...")
        ppt_path = create_enhanced_ppt(content, reference or "", username)

        return FileResponse(
            ppt_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="enhanced_from_doc.pptx"
        )

    except Exception as e:
        print(f"Unexpected error: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")
    finally:
        if os.path.exists(path):
            try:
                os.remove(path)
            except Exception as e:
                print(f"Cleanup error: {e}")


@app.on_event("startup")
async def startup_event():
    """Test summarizer functionality on startup"""
    print("="*50)
    print("TESTING SUMMARIZER FUNCTIONALITY")
    print("="*50)
    
    test_text = """
    Artificial intelligence is transforming the world. Machine learning algorithms can process vast amounts of data.
    Natural language processing enables computers to understand human language. Deep learning networks mimic the human brain.
    These technologies are revolutionizing industries from healthcare to finance.
    """
    
    try:
        
        title = summarizer.generate_title(test_text)
        bullets = summarizer.generate_bullet_points(test_text, max_bullets=3)
        
        print(f"✓ Title generation: '{title}'")
        print(f"✓ Bullet generation: {len(bullets)} bullets created")
        for i, bullet in enumerate(bullets, 1):
            print(f"  {i}. {bullet}")
        
        print("✓ Summarizer is working correctly!")
    except Exception as e:
        print(f"✗ Summarizer test failed: {e}")
        import traceback
        traceback.print_exc()
    
    print("="*50)

@app.get("/test-summarizer")
async def test_summarizer_endpoint():
    """Test endpoint for summarizer functionality"""
    test_text = """
    Machine learning is a subset of artificial intelligence that enables computers to learn and improve from experience.
    It involves algorithms that can identify patterns in data and make predictions or decisions.
    Applications include image recognition, natural language processing, and recommendation systems.
    The field continues to evolve with new techniques like deep learning and neural networks.
    """
    
    try:
        title = summarizer.generate_title(test_text)
        bullets = summarizer.generate_bullet_points(test_text, max_bullets=3)
        
        return JSONResponse({
            "status": "success",
            "title": title,
            "bullets": bullets,
            "nltk_available": getattr(summarizer, 'nltk_available', False)
        })
    except Exception as e:
        return JSONResponse({
            "status": "error",
            "error": str(e),
            "nltk_available": getattr(summarizer, 'nltk_available', False)
        })






@app.get("/user-history/{username}")
async def get_user_history(username: str):
    logs = load_logs()
    return JSONResponse(logs.get(username, {"browser": "unknown", "history": []}))
